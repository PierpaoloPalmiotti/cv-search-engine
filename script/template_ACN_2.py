from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import json, re

# ---------- Utils ----------
TAG = re.compile(r"\{\{\s*([^\}]+?)\s*\}\}", re.IGNORECASE)

def has_tag(shape, tag_norm):
    """Verifica se uno shape contiene un tag specifico"""
    if not shape.has_text_frame:
        return False
    full = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
    norm = TAG.sub(lambda m: "{{" + m.group(1).strip().upper() + "}}", full)
    return tag_norm in norm

def get_all_tags(shape):
    """Estrae tutti i tag presenti in uno shape"""
    if not shape.has_text_frame:
        return []
    full = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
    matches = TAG.findall(full)
    return [m.strip().upper() for m in matches]

def replace_tag_in_place(shape, tag_name, replacement_text, font_size=None, bold=None):
    """Sostituisce un tag preservando il resto del contenuto"""
    tag_pattern = re.compile(r"\{\{\s*" + re.escape(tag_name) + r"\s*\}\}", re.IGNORECASE)
    
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            if tag_pattern.search(run.text):
                run.text = tag_pattern.sub(replacement_text, run.text)
                if font_size:
                    run.font.size = Pt(font_size)
                if bold is not None:
                    run.font.bold = bold
                run.font.name = "Arial"

def set_text(shape, text, font_size=None, bold=None, align=None, line_spacing=1.0):
    """Sostituisce completamente il contenuto di uno shape"""
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align is not None:
            p.alignment = align
        for r in p.runs:
            if font_size:
                r.font.size = Pt(font_size)
            if bold is not None:
                r.font.bold = bold
            r.font.name = "Arial"
        p.line_spacing = line_spacing

def fill_bullets(shape, items, limit=6, font_size=11):
    """Riempie uno shape con bullet points"""
    items = (items or [])[:limit]
    tf = shape.text_frame
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        p.space_after = Pt(2)

def wrap_title(title):
    """Formatta il titolo con a capo intelligenti"""
    t = re.sub(r"\s*/\s*", " / ", title.strip())
    t = re.sub(r",\s*", ",\n", t)
    if len(title) > 45 and "\n" not in t:
        t = t.replace(" & ", "\n& ").replace(" and ", "\nand ")
    return t

def autosize_for_title(text):
    """Calcola dimensione font ottimale per il titolo"""
    ln = len(text.replace("\n"," "))
    if ln <= 28: return 32
    if ln <= 45: return 28
    if ln <= 60: return 24
    if ln <= 80: return 20
    return 18

# ---------- FUNZIONI PER PROJECT EXPERIENCE ----------

def add_para(tf, text, size=10, bold=False, before=0, after=0):
    """Aggiunge un paragrafo senza bullet"""
    p = tf.add_paragraph() if len(tf.paragraphs) and tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = ""
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = "Arial"
    p.level = 0
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.alignment = PP_ALIGN.LEFT
    return p

def add_label_value(tf, label, value, size=10, before=0, after=2):
    """Crea una riga con label bold + valore normale"""
    p = tf.add_paragraph()
    p.text = ""
    # Label in grassetto
    r1 = p.add_run()
    r1.text = label
    r1.font.bold = True
    r1.font.size = Pt(size)
    r1.font.name = "Arial"
    # Valore normale
    r2 = p.add_run()
    r2.text = value
    r2.font.bold = False
    r2.font.size = Pt(size)
    r2.font.name = "Arial"
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    p.alignment = PP_ALIGN.LEFT
    return p

def fill_project_experience(shape, experiences, max_items=5):
    """Riempie il box delle esperienze professionali"""
    tf = shape.text_frame
    
    # Azzera margini
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    
    tf.clear()
    
    # Assicurati che esista almeno un paragrafo
    if not tf.paragraphs:
        tf.paragraphs.add()
    
    for idx, exp in enumerate(experiences[:max_items]):
        company = (exp.get("company","") or "").upper().strip()
        location = exp.get("location","")
        period = exp.get("period","")
        
        # Header dell'esperienza
        head = f"{company}"
        if location:
            head += f" ({location})"
        if period:
            head += f"  {period}"

        # Prima esperienza: usa il primo paragrafo esistente con spacing zero
        if idx == 0:
            p = tf.paragraphs[0]
            p.text = ""
            r = p.add_run()
            r.text = head
            r.font.size = Pt(11)
            r.font.bold = True
            r.font.name = "Arial"
            p.level = 0
            p.space_before = Pt(0)
            p.space_after = Pt(1)
            p.alignment = PP_ALIGN.LEFT
        else:
            add_para(tf, head, size=11, bold=True, before=6, after=1)

        # Responsibilities
        resp = exp.get("responsibilities") or exp.get("description") or ""
        if isinstance(resp, list):
            resp = " • ".join([s.strip() for s in resp if s.strip()])
        if resp:
            add_label_value(tf, "Responsibilities: ", resp, size=10, before=0, after=0.5)

        # Technologies
        techs = exp.get("technologies", [])
        if isinstance(techs, list):
            techs = ", ".join([t for t in techs if t])
        if techs:
            add_label_value(tf, "Technologies: ", techs, size=10, before=0, after=2)

# ---------- MAPPING TAG UNIVERSALE ----------
# Definisci qui le possibili varianti di tag per ogni campo
TAG_MAPPING = {
    "name": ["NOME", "NOME COMPLETO", "NOME_CV", "NAME"],
    "title": ["TITOLO", "TITLE", "RUOLO", "ROLE"],
    "background": ["BACKGROUND", "RIASSUNTO", "SUMMARY", "PROFESSIONALE"],
    "skills": ["COMPETENZE", "COMPETENZA", "SKILLS", "SKILL"],
    "technologies": ["TECNOLOGIE", "TECNOLOGIA", "TECH", "TECHNOLOGIES"],
    "experience": ["ESPERIENZE", "ESPERIENZA", "EXPERIENCE", "PROJECT EXPERIENCE", "PROFESSIONALE"]
}

def find_tag_match(shape, category):
    """Trova quale variante di tag è presente nello shape"""
    shape_tags = get_all_tags(shape)
    for tag_variant in TAG_MAPPING.get(category, []):
        if tag_variant in shape_tags:
            return tag_variant
    return None

# ---------- CARICAMENTO DATI ----------
with open("info2.json", "r", encoding="utf-8") as f:
    data = json.load(f)

name = data.get("name","").strip()
title_raw = data.get("title","").strip()
background = data.get("summary","").strip()
skills = data.get("skills", [])
techs = data.get("technologies", [])

# Prepara esperienze
raw_exp = data.get("experience", [])
experiences = []
for e in raw_exp:
    experiences.append({
        "company": e.get("company",""),
        "location": e.get("location",""),
        "period": e.get("period",""),
        "responsibilities": e.get("responsibilities") or e.get("description",""),
        "technologies": e.get("technologies", e.get("tech", [])),
    })

# ---------- PROCESSAMENTO TEMPLATE ----------
prs = Presentation("CV_template_standardizzato2.pptx")

# Debug: stampa tutti i tag trovati
print("\n=== DEBUG: Tag trovati nel template ===")
for idx, slide in enumerate(prs.slides):
    for sh_idx, sh in enumerate(slide.shapes):
        if sh.has_text_frame:
            tags = get_all_tags(sh)
            if tags:
                print(f"Slide {idx}, Shape {sh_idx}: {tags}")
print("========================================\n")

for slide in prs.slides:
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue

        # Ottieni tutti i tag presenti nello shape
        all_tags = get_all_tags(sh)
        
        # Se lo shape contiene sia NOME che TITOLO, gestiscili insieme
        name_tag = find_tag_match(sh, "name")
        title_tag = find_tag_match(sh, "title")
        
        if name_tag and title_tag:
            print(f"✓ Compilando NOME e TITOLO insieme")
            # Sostituisci entrambi i tag preservando la struttura
            replace_tag_in_place(sh, name_tag, name, font_size=44, bold=True)
            t = wrap_title(title_raw)
            # Font più piccolo per il titolo per evitare sovrapposizioni
            replace_tag_in_place(sh, title_tag, t, font_size=18, bold=True)
            continue
        
        # NOME (solo se non c'è anche TITOLO nello stesso shape)
        if name_tag:
            print(f"✓ Compilando NOME con: {name}")
            set_text(sh, name, font_size=44, bold=True, align=PP_ALIGN.LEFT, line_spacing=1.0)
            continue

        # TITOLO (solo se non c'è anche NOME nello stesso shape)
        if title_tag:
            t = wrap_title(title_raw)
            print(f"✓ Compilando TITOLO con: {t}")
            set_text(sh, t, font_size=autosize_for_title(t), bold=True, align=PP_ALIGN.LEFT, line_spacing=1.0)
            continue

        # BACKGROUND - prova prima sostituzione in place
        tag = find_tag_match(sh, "background")
        if tag:
            print(f"✓ Compilando BACKGROUND")
            # Verifica se ci sono altri testi oltre al tag
            full_text = "".join(run.text for p in sh.text_frame.paragraphs for run in p.runs)
            tag_text = f"{{{{{tag}}}}}"
            
            # Se c'è solo il tag, sostituisci tutto
            if full_text.strip().replace(tag_text, "").strip() == "":
                set_text(sh, background, font_size=9, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.0)
            else:
                # Altrimenti sostituisci solo il tag
                replace_tag_in_place(sh, tag, background, font_size=9, bold=False)
            continue

        # COMPETENZE e TECNOLOGIE - gestisci se sono nello stesso shape
        skills_tag = find_tag_match(sh, "skills")
        tech_tag = find_tag_match(sh, "technologies")
        
        if skills_tag and tech_tag:
            print(f"✓ Compilando COMPETENZE e TECNOLOGIE insieme")
            # Se sono insieme, sostituisci i tag separatamente
            replace_tag_in_place(sh, skills_tag, "\n".join(f"• {s}" for s in skills[:6]), font_size=11, bold=False)
            replace_tag_in_place(sh, tech_tag, "\n".join(f"• {t}" for t in techs[:8]), font_size=11, bold=False)
            continue
        
        # COMPETENZE (solo se non c'è anche TECNOLOGIE)
        if skills_tag:
            print(f"✓ Compilando COMPETENZE ({len(skills)} items)")
            fill_bullets(sh, skills, limit=6, font_size=11)
            continue

        # TECNOLOGIE (solo se non c'è anche COMPETENZE)
        if tech_tag:
            print(f"✓ Compilando TECNOLOGIE ({len(techs)} items)")
            fill_bullets(sh, techs, limit=8, font_size=11)
            continue

        # ESPERIENZE
        tag = find_tag_match(sh, "experience")
        if tag:
            print(f"✓ Compilando ESPERIENZE ({len(experiences)} items)")
            fill_project_experience(sh, experiences, max_items=5)
            continue

# Salva
prs.save("cv_template2_compilato.pptx")
print("✓ CV generato: cv_template2_compilato.pptx")