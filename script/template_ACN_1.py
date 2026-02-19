from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import json
import re

# --- 1. Leggi JSON ---
with open("info.json", "r", encoding="utf-8") as f:
    data = json.load(f)

# --- 2. Mappa base ---
campi = {
    "{{NOME}}": data.get("name", ""),
    "{{TITOLO}}": data.get("title", ""),
    # "{{RIASSUNTO}}": data.get("summary", ""),  # Gestito separatamente
    "{{FORMAZIONE}}": data['education'].get('degree', ''),
    "{{ANNOFINE_FORMAZIONE}}": str(data["education"].get("year", "")) if data["education"].get("year") else "",
    "{{DATAINIZIOCERTI}}": "",  # Rimosso, sostituito con stringa vuota
    "{{DATAFINECERT}}": "",      # Rimosso, sostituito con stringa vuota
}

# --- 3. Funzione di sostituzione base (per campi singoli) ---
def replace_text_in_shape(shape, replacements):
    if not shape.has_text_frame:
        return
    
    # Itera sui paragrafi e run per mantenere la formattazione
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for tag, val in replacements.items():
                if tag in run.text:
                    run.text = run.text.replace(tag, val)
                    # Rendi il nome in grassetto
                    if tag == "{{NOME}}":
                        run.font.bold = True
            
            # Rimuovi trattini isolati che rimangono dopo placeholder vuoti
            run.text = run.text.replace(" - ", " ").replace(" -", "").strip()

# --- 3b. Funzione per sostituire testo con dimensione font specifica ---
def replace_text_with_font_size(shape, tag, value, font_size):
    """Sostituisce il testo e imposta una dimensione font specifica"""
    if not shape.has_text_frame:
        return
    
    full_text = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
    
    if tag in full_text:
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.size = Pt(font_size)
        run.font.name = "Arial"

# --- 4. Carica presentazione ---
prs = Presentation("CV_template_standardizzato2.pptx")

# --- 5. Sostituisci testo semplice ---
for slide in prs.slides:
    for shape in slide.shapes:
        replace_text_in_shape(shape, campi)

# --- 5b. Riduci font del riassunto ---
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame and "{{RIASSUNTO}}" in shape.text:
            replace_text_with_font_size(shape, "{{RIASSUNTO}}", data.get("summary", ""), 9)

# --- 6. Inserisci elenchi formattati ---
def fill_list(shape, items):
    """Sostituisce il testo del shape con un elenco puntato"""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    
    for i, item in enumerate(items):
        # Usa il primo paragrafo esistente, poi aggiungi i successivi
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item
        p.level = 0
        p.space_after = Pt(4)
        p.space_before = Pt(4)
        p.font.size = Pt(11)
        p.font.name = "Arial"
        p.alignment = PP_ALIGN.LEFT
        p.font.bold = False

# --- 7. Inserisci sezioni multiple con formato simile al modello ---
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text
        # COMPETENZE
        if "{{COMPETENZA}}" in shape.text:
            fill_list(shape, data.get("skills", []))
        # TECNOLOGIE
        elif "{{TECNOLOGIE}}" in shape.text:
            fill_list(shape, data.get("technologies", []))
        # ESPERIENZE
        elif "{{ESPERIENZE}}" in shape.text:
            esperienze = [
                f"{exp['company'].upper()} ({exp['period']}):\n{exp['description']}"
                for exp in data.get("experience", [])
            ]
            fill_list(shape, esperienze)
        # CERTIFICAZIONI - Non usare fill_list, è già gestito dai placeholder individuali
        elif "{{CERTIFICAZIONE}}" in shape.text:
            # Prendi la prima certificazione
            certifications = data.get("certifications", [])
            if certifications:
                cert_name = certifications[0] if isinstance(certifications[0], str) else certifications[0].get("name", "")
                # Sostituisci solo il placeholder, non usare fill_list
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace("{{CERTIFICAZIONE}}", cert_name)

# --- 7b. Uniforma dimensione font nella sezione ISTRUZIONE & FORMAZIONE ---
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Cerca il titolo della sezione per identificarla
        full_text = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
        if "ISTRUZIONE" in full_text and "FORMAZIONE" in full_text:
            # Salta il titolo della sezione, uniforma solo il contenuto sotto
            continue
        # Uniforma i testi che contengono i dati della formazione
        if any(keyword in full_text for keyword in ["Laurea", "Google Cloud", "Microsoft", "{{FORMAZIONE}}"]):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Non modificare se è il titolo (grassetto grande)
                    if run.font.size and run.font.size > Pt(12):
                        continue
                    run.font.size = Pt(11)
                    run.font.name = "Arial"

# --- 8. Salva file ---
prs.save("veneranda_taldone_format.pptx")
print("File 'veneranda_taldone_format.pptx' generato con formattazione migliorata!")