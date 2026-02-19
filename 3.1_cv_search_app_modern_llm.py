# cv_search_app_modern.py
"""
INSTALLAZIONE RICHIESTA:
pip install customtkinter
"""

import customtkinter as ctk
import threading
import matplotlib.pyplot as plt
from mpl_toolkits.mplot3d import Axes3D
from sklearn.decomposition import PCA
from tkinter import messagebox
import numpy as np
from FlagEmbedding import BGEM3FlagModel
from sklearn.metrics.pairwise import cosine_similarity
import os
import json
import re
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import requests


# Configura tema e colori
ctk.set_appearance_mode("dark")  # "dark" o "light"
ctk.set_default_color_theme("blue")  # "blue", "green", "dark-blue"

class Logger:
    """Gestisce il logging su file"""
    def __init__(self, log_file="cv_search_log.txt"):
        self.log_file = log_file
        self.log_to_file(f"\n{'='*80}\nNuova sessione iniziata: {datetime.now()}\n{'='*80}\n")
    
    def log_to_file(self, message):
        """Scrive sul file di log"""
        with open(self.log_file, 'a', encoding='utf-8') as f:
            f.write(f"{message}\n")
    
    def log(self, message, level="INFO"):
        """Log con timestamp"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        log_msg = f"[{timestamp}] [{level}] {message}"
        self.log_to_file(log_msg)
        print(log_msg)

class PPTXToJSONExtractor:
    """Estrattore PPTX -> JSON"""
    def __init__(self, cv_ppt_folder="./cv_ppt", cv_json_folder="./cv_json", logger=None):
        self.cv_ppt_folder = Path(cv_ppt_folder)
        self.cv_json_folder = Path(cv_json_folder)
        self.cv_json_folder.mkdir(exist_ok=True, parents=True)
        self.logger = logger or Logger()
    
    def extract_text_from_shapes(self, slide):
        """Estrae tutto il testo dalle shape della slide"""
        all_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
        return "\n".join(all_text)
    
    def find_pptx_by_label(self, label):
        """Trova il file PPTX corrispondente al label"""
        normalized = label.lower().strip().replace(' ', '_').replace('.', '_')
        
        for pptx_file in self.cv_ppt_folder.glob("*.pptx"):
            filename_lower = pptx_file.stem.lower()
            if normalized in filename_lower or filename_lower in normalized:
                return pptx_file
        
        self.logger.log(f"PPTX non trovato per: {label}", "WARNING")
        return None
    
    def extract_name(self, text):
        """Estrae il nome completo"""
        lines = text.split('\n')
        patterns = [
            r'^([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)',
            r'(?:Nome|Name)[:\s]+([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)',
        ]
        for line in lines[:5]:
            for pattern in patterns:
                match = re.search(pattern, line)
                if match:
                    name = match.group(1).strip()
                    if len(name.split()) >= 2:
                        return name.title()
        return ""
    
    def extract_title(self, text):
        """Estrae il titolo/ruolo professionale"""
        keywords = ['ruolo', 'posizione', 'title', 'role']
        lines = text.split('\n')
        for i, line in enumerate(lines):
            if any(kw in line.lower() for kw in keywords):
                title = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if title:
                    return title
                if i + 1 < len(lines):
                    return lines[i + 1].strip()
        return ""
    
    def extract_summary(self, text):
        """Estrae il sommario/profilo"""
        keywords = ['profilo', 'summary', 'about', 'descrizione']
        lines = text.split('\n')
        summary_lines = []
        capturing = False
        
        for line in lines:
            if any(kw in line.lower() for kw in keywords):
                capturing = True
                content = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if content and len(content) > 20:
                    summary_lines.append(content)
                continue
            
            if capturing:
                if any(stop in line.lower() for stop in ['competenze', 'formazione', 'esperienza']):
                    break
                if line.strip() and len(line.strip()) > 10:
                    summary_lines.append(line.strip())
                if len(summary_lines) >= 5:
                    break
        
        return " ".join(summary_lines)
    
    def extract_list_items(self, text, keywords):
        """Estrae lista di elementi"""
        lines = text.split('\n')
        items = []
        capturing = False
        
        for line in lines:
            if any(kw in line.lower() for kw in keywords):
                capturing = True
                continue
            
            if capturing:
                if any(stop in line.lower() for stop in ['formazione', 'esperienza', 'certificazioni']):
                    break
                if line.strip():
                    parsed = self._parse_items(line)
                    items.extend(parsed)
        
        return list(dict.fromkeys(items))
    
    def _parse_items(self, text):
        """Parse una riga in items individuali"""
        text = re.sub(r'^[\-\•\*\d\.\)]+\s*', '', text.strip())
        for sep in [',', ';', '|']:
            if sep in text:
                return [item.strip() for item in text.split(sep) if len(item.strip()) > 1]
        return [text] if len(text) > 1 else []
    
    def extract_experience(self, text):
        """Estrae esperienze lavorative"""
        keywords = ['esperienza', 'experience', 'lavoro']
        lines = text.split('\n')
        experiences = []
        capturing = False
        current_exp = []
        
        for line in lines:
            if any(kw in line.lower() for kw in keywords):
                capturing = True
                continue
            
            if capturing:
                if any(stop in line.lower() for stop in ['formazione', 'competenze']):
                    break
                
                if line.strip():
                    if re.search(r'\b(19|20)\d{2}\b', line):
                        if current_exp:
                            experiences.append(self._parse_experience(current_exp))
                            current_exp = []
                    current_exp.append(line.strip())
        
        if current_exp:
            experiences.append(self._parse_experience(current_exp))
        
        return experiences
    
    def _parse_experience(self, lines):
        """Parse singola esperienza"""
        full_text = " ".join(lines)
        experience = {
            "company": lines[0] if lines else "",
            "period": "",
            "description": " ".join(lines[1:]) if len(lines) > 1 else "",
            "location": "",
            "responsibilities": " ".join(lines[1:]) if len(lines) > 1 else "",
            "technologies": []
        }
        
        period_match = re.search(r'(\d{4}\s*-\s*\d{4})', full_text)
        if period_match:
            experience["period"] = period_match.group(1)
        
        return experience
    
    def extract_info_from_pptx(self, pptx_path):
        """Estrae tutte le informazioni da un file PPTX"""
        try:
            prs = Presentation(str(pptx_path))
            if len(prs.slides) == 0:
                self.logger.log(f"Nessuna slide in {pptx_path.name}", "WARNING")
                return self._empty_json()
            
            slide = prs.slides[0]
            all_text = self.extract_text_from_shapes(slide)
            
            if not all_text.strip():
                self.logger.log(f"Nessun testo estratto da {pptx_path.name}", "WARNING")
                return self._empty_json()
            
            json_data = {
                "name": self.extract_name(all_text),
                "title": self.extract_title(all_text),
                "summary": self.extract_summary(all_text),
                "skills": self.extract_list_items(all_text, ['competenze', 'skills']),
                "technologies": self.extract_list_items(all_text, ['tecnologie', 'technologies']),
                "certifications": self.extract_list_items(all_text, ['certificazioni', 'certifications']),
                "experience": self.extract_experience(all_text)
            }
            
            return json_data
            
        except Exception as e:
            self.logger.log(f"Errore processando {pptx_path.name}: {e}", "ERROR")
            return self._empty_json()
    
    def _empty_json(self):
        """JSON vuoto"""
        return {
            "name": "", "title": "", "summary": "",
            "skills": [], "technologies": [], "certifications": [], "experience": []
        }
    
    def process_label(self, label):
        """Processa un singolo label: PPTX -> JSON"""
        self.logger.log(f"Tentativo estrazione da PPTX per: {label}")
        
        pptx_file = self.find_pptx_by_label(label)
        if not pptx_file:
            self.logger.log(f"PPTX non trovato per: {label}", "WARNING")
            return None
        
        self.logger.log(f"PPTX trovato: {pptx_file.name}")
        json_data = self.extract_info_from_pptx(pptx_file)
        
        json_filename = self.cv_json_folder / f"{pptx_file.stem}.json"
        with open(json_filename, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
        
        self.logger.log(f"JSON estratto e salvato: {json_filename.name}")
        return json_filename

class PPTXGeneratorACN1:
    """Generatore PPTX da JSON usando template ACN_1"""
    def __init__(self, template_path, logger=None):
        self.template_path = Path(template_path)
        self.logger = logger or Logger()
        self.template_name = "ACN_1"
    
    def replace_text_in_shape(self, shape, replacements):
        """Sostituisce testo nei run mantenendo la formattazione"""
        if not shape.has_text_frame:
            return
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for tag, val in replacements.items():
                    if tag in run.text:
                        run.text = run.text.replace(tag, val)
                        if tag == "{{NOME}}":
                            run.font.bold = True
                run.text = run.text.replace(" - ", " ").replace(" -", "").strip()
    
    def replace_text_with_font_size(self, shape, tag, value, font_size):
        """Sostituisce il testo e imposta una dimensione font specifica"""
        if not shape.has_text_frame:
            return
        
        full_text = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
        
        if tag in full_text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = value
            run.font.size = Pt(font_size)
            run.font.name = "Arial"
    
    def fill_list(self, shape, items):
        """Sostituisce il testo del shape con un elenco puntato"""
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        tf.clear()
        
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = item
            p.level = 0
            p.space_after = Pt(4)
            p.space_before = Pt(4)
            p.font.size = Pt(11)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.LEFT
            p.font.bold = False
        
    def generate_cv(self, json_path, output_path):
        """Genera CV da JSON usando template ACN_1"""
        self.logger.log(f"Generando CV con template {self.template_name} da: {json_path}")
        
        if not self.template_path.exists():
            self.logger.log(f"Template non trovato: {self.template_path}", "ERROR")
            return False
        
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            campi = {
                "{{NOME}}": data.get("name", ""),
                "{{TITOLO}}": data.get("title", ""),
                "{{FORMAZIONE}}": data.get('education', {}).get('degree', ''),
                "{{ANNOFINE_FORMAZIONE}}": str(data.get("education", {}).get("year", "")) if data.get("education", {}).get("year") else "",
                "{{DATAINIZIOCERTI}}": "",
                "{{DATAFINECERT}}": "",
            }
            
            prs = Presentation(str(self.template_path))
            
            # STEP 1: Sostituzioni semplici
            for slide in prs.slides:
                for shape in slide.shapes:
                    self.replace_text_in_shape(shape, campi)
            
            # STEP 2: Gestione BACKGROUND (era RIASSUNTO)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and "{{BACKGROUND}}" in shape.text:
                        self.replace_text_with_font_size(shape, "{{BACKGROUND}}", data.get("summary", ""), 9)
            
            # STEP 3: Gestione liste complesse
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    text = shape.text
                    
                    # SKILLS unificato (era COMPETENZA + TECNOLOGIE)
                    if "{{SKILLS}}" in text:
                        all_skills = data.get("skills", []) + data.get("technologies", [])
                        self.fill_list(shape, all_skills)
                    
                    # ESPERIENZE
                    elif "{{ESPERIENZE}}" in text:
                        esperienze = [
                            f"{exp.get('company', '').upper()} ({exp.get('period', '')}):\n{exp.get('description', '')}"
                            for exp in data.get("experience", [])
                        ]
                        self.fill_list(shape, esperienze)
                    
                    # CERTIFICAZIONI
                    elif "{{CERTIFICAZIONI}}" in text:
                        certifications = data.get("certifications", [])
                        if certifications:
                            cert_list = []
                            for cert in certifications:
                                if isinstance(cert, str):
                                    cert_list.append(cert)
                                else:
                                    cert_list.append(cert.get("name", ""))
                            self.fill_list(shape, cert_list)
            
            # STEP 4: Formattazione font (opzionale)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    full_text = "".join(run.text for p in shape.text_frame.paragraphs for run in p.runs)
                    
                    if "ISTRUZIONE" in full_text and "FORMAZIONE" in full_text:
                        continue
                    
                    if any(keyword in full_text for keyword in ["Laurea", "Google Cloud", "Microsoft", "{{FORMAZIONE}}"]):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size and run.font.size > Pt(12):
                                    continue
                                run.font.size = Pt(11)
                                run.font.name = "Arial"
            
            # SALVA
            prs.save(str(output_path))
            self.logger.log(f"CV generato con successo: {output_path}")
            return True
            
        except Exception as e:
            self.logger.log(f"Errore generazione CV: {e}", "ERROR")
            return False


class PPTXGeneratorGeneric:
    """Generatore PPTX generico per template con placeholder standard"""
    def __init__(self, template_path, logger=None):
        self.template_path = Path(template_path)
        self.logger = logger or Logger()
        self.template_name = template_path.stem
    
    def replace_text_in_shape(self, shape, replacements):
        """Sostituisce testo preservando formattazione"""
        if not shape.has_text_frame:
            return
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for tag, val in replacements.items():
                    if tag in run.text:
                        run.text = run.text.replace(tag, str(val))
    
    def fill_list(self, shape, items, as_bullet=True):
        """Riempie con lista di items"""
        if not shape.has_text_frame:
            return
        
        tf = shape.text_frame
        tf.clear()
        
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            if as_bullet:
                p.level = 0
            p.font.size = Pt(11)
            p.font.name = "Arial"
    
    def generate_cv(self, json_path, output_path):
        """Genera CV da JSON usando template generico"""
        self.logger.log(f"Generando CV con template {self.template_name} da: {json_path}")
        
        if not self.template_path.exists():
            self.logger.log(f"Template non trovato: {self.template_path}", "ERROR")
            return False
        
        try:
            with open(json_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            # ── Mapping unificato (aggiunto {{RUOLO}} come alias di TITOLO) ──
            campi = {
                "{{NOME}}": data.get("name", ""),
                "{{TITOLO}}": data.get("title", ""),
                "{{RUOLO}}": data.get("title", ""),          # ← NUOVO
                "{{BACKGROUND}}": data.get("summary", ""),
                "{{FORMAZIONE}}": data.get('education', {}).get('degree', ''),
            }
            
            prs = Presentation(str(self.template_path))
            
            # STEP 1: Sostituzioni semplici
            for slide in prs.slides:
                for shape in slide.shapes:
                    self.replace_text_in_shape(shape, campi)
            
            # ── Prepara split delle skills ──────────────────────────────
            all_skills = data.get("skills", []) + data.get("technologies", [])
            
            if len(all_skills) > 5:
                skills_col1 = all_skills[:5]
                skills_col2 = all_skills[5:]
            else:
                skills_col1 = all_skills
                skills_col2 = []          # vuoto → verrà scritto blank
            
            self.logger.log(
                f"Skills split: col1={len(skills_col1)}, col2={len(skills_col2)} "
                f"(totale {len(all_skills)})"
            )
            # ────────────────────────────────────────────────────────────
            
            # STEP 2: Gestione liste e sezioni complesse
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    text = shape.text
                    
                    # ESPERIENZE
                    if "{{ESPERIENZE}}" in text:
                        experiences_text = []
                        for exp in data.get("experience", []):
                            exp_line = (
                                f"• {exp.get('company', '')} "
                                f"({exp.get('period', '')}): "
                                f"{exp.get('description', '')}"
                            )
                            experiences_text.append(exp_line)
                        
                        if experiences_text:
                            shape.text_frame.clear()
                            for i, exp_text in enumerate(experiences_text):
                                p = (shape.text_frame.paragraphs[0]
                                     if i == 0
                                     else shape.text_frame.add_paragraph())
                                p.text = exp_text
                                p.font.size = Pt(10)
                                p.font.name = "Arial"
                        self.logger.log(f"Compilato ESPERIENZE: {len(experiences_text)} items")
                    
                    # ── SKILLS SPLIT: colonna 1 ────────────────────────
                    elif "{{SKILLS1}}" in text:
                        self.fill_list(shape, skills_col1)
                        self.logger.log(f"Compilato SKILLS1: {len(skills_col1)} items")
                    
                    # ── SKILLS SPLIT: colonna 2 ────────────────────────
                    elif "{{SKILLS2}}" in text:
                        if skills_col2:
                            self.fill_list(shape, skills_col2)
                            self.logger.log(f"Compilato SKILLS2: {len(skills_col2)} items")
                        else:
                            # Nessuna skill residua → scrivi blank
                            shape.text_frame.clear()
                            p = shape.text_frame.paragraphs[0]
                            p.text = ""
                            self.logger.log("SKILLS2: blank (≤ 5 skills totali)")
                    
                    # ── Retrocompatibilità: tag singolo {{SKILLS}} ─────
                    elif "{{SKILLS}}" in text:
                        self.fill_list(shape, all_skills)
                        self.logger.log(f"Compilato SKILLS (singolo): {len(all_skills)} items")
                    
                    # CERTIFICAZIONI
                    elif "{{CERTIFICAZIONI}}" in text:
                        certifications = data.get("certifications", [])
                        cert_list = []
                        for cert in certifications:
                            if isinstance(cert, str):
                                cert_list.append(cert)
                            else:
                                cert_list.append(cert.get("name", ""))
                        self.fill_list(shape, cert_list)
                        self.logger.log(f"Compilato CERTIFICAZIONI: {len(cert_list)} items")
                    
                    # LINGUE
                    elif "{{LINGUE}}" in text:
                        lingue = data.get("languages",
                                          ["Italiano (madrelingua)", "Inglese (fluente)"])
                        self.fill_list(shape, lingue)
                        self.logger.log(f"Compilato LINGUE: {len(lingue)} items")
            
            # SALVA
            prs.save(str(output_path))
            self.logger.log(f"CV generato con successo: {output_path}")
            return True
            
        except Exception as e:
            self.logger.log(f"Errore generazione CV con {self.template_name}: {e}", "ERROR")
            return False


def get_available_templates(template_folder="./template"):
    """Scansiona la cartella template e restituisce i template disponibili"""
    template_path = Path(template_folder)
    if not template_path.exists():
        return []
    
    templates = []
    for pptx_file in template_path.glob("*.pptx"):
        if not pptx_file.name.startswith("~$"):  # Ignora file temporanei
            templates.append({
                "name": pptx_file.stem,
                "path": pptx_file,
                "display_name": pptx_file.stem.replace("_", " ").title()
            })
    
    return sorted(templates, key=lambda x: x["name"])


def create_generator_for_template(template_path, logger):
    """Usa sempre il generator generico (tag unificati)"""
    return PPTXGeneratorGeneric(template_path, logger)

class CVSearchApp:
    def __init__(self):
        self.root = ctk.CTk()
        self.root.title("CV Search Engine - Modern UI")
        self.root.geometry("1200x850")
        
        # Logger
        self.logger = Logger()
        
        # Variabili
        self.model = None
        self.cv_embeddings = None
        self.cv_texts = None
        self.cv_labels = None
        self.selected_template = None
        self.available_templates = []
        self.selected_llm_model = "llama3.2:1b"  # Modello di default
        
        # Setup UI
        self.setup_ui()
        self.load_templates()
        self.load_data()
    
    def setup_ui(self):
        # Status bar compatto in alto
        status_frame = ctk.CTkFrame(self.root, corner_radius=0, fg_color=("gray90", "gray13"), height=35)
        status_frame.pack(fill="x", padx=0, pady=0)
        status_frame.pack_propagate(False)
        
        self.status_label = ctk.CTkLabel(status_frame, text="⏳ Caricamento in corso...",
                                        font=ctk.CTkFont(size=11), text_color="orange")
        self.status_label.pack(pady=8)
        
        # Main Content Frame con layout a 2 colonne
        content_frame = ctk.CTkFrame(self.root, fg_color="transparent")
        content_frame.pack(fill="both", expand=True, padx=20, pady=(0, 20))
        
        # Configura griglia per 2 colonne di uguale larghezza
        content_frame.grid_columnconfigure(0, weight=1, uniform="col")
        content_frame.grid_columnconfigure(1, weight=1, uniform="col")
        content_frame.grid_rowconfigure(0, weight=1)
        
        # ===== COLONNA SINISTRA: Query e Controlli =====
        left_column = ctk.CTkFrame(content_frame, corner_radius=10)
        left_column.grid(row=0, column=0, sticky="nsew", padx=(0, 5))
        
        # SEZIONE GENERAZIONE DIRETTA
        direct_gen_header = ctk.CTkFrame(left_column, fg_color=("gray85", "gray17"), corner_radius=8)
        direct_gen_header.pack(fill="x", padx=10, pady=(0, 5))
        
        direct_gen_label = ctk.CTkLabel(direct_gen_header, text="👤 Generazione Diretta CV",
                                        font=ctk.CTkFont(size=14, weight="bold"))
        direct_gen_label.pack(pady=8)
        
        direct_gen_frame = ctk.CTkFrame(left_column, fg_color=("gray90", "gray20"), corner_radius=8)
        direct_gen_frame.pack(fill="x", padx=10, pady=(0, 10))
        
        ctk.CTkLabel(direct_gen_frame, text="Nome e Cognome:", 
                    font=ctk.CTkFont(size=11, weight="bold")).pack(anchor="w", padx=10, pady=(10, 2))
        
        self.direct_name_entry = ctk.CTkEntry(direct_gen_frame,
                                              placeholder_text="Es: Mario Rossi",
                                              font=ctk.CTkFont(size=11),
                                              height=32)
        self.direct_name_entry.pack(fill="x", padx=10, pady=(0, 10))
        
        self.direct_gen_button = ctk.CTkButton(direct_gen_frame,
                                               text="⚡ Genera CV Diretto",
                                               command=self.generate_direct_cv,
                                               font=ctk.CTkFont(size=12, weight="bold"),
                                               height=38,
                                               corner_radius=8,
                                               fg_color=("#2E7D32", "#1B5E20"),
                                               hover_color=("#1B5E20", "#0D3F10"))
        self.direct_gen_button.pack(fill="x", padx=10, pady=(0, 10))
        
        self.query_text = ctk.CTkTextbox(left_column, 
                                        font=ctk.CTkFont(size=11),
                                        corner_radius=8,
                                        wrap="word",
                                        height=100)
        self.query_text.pack(fill="x", padx=10, pady=(0, 10))
        self.query_text.insert("1.0", "Skills: ...\nIndustry: ...\nOffice: ...\n Level: ... ")
       
        control_content = ctk.CTkFrame(left_column, fg_color="transparent")
        control_content.pack(fill="x", padx=10, pady=(0, 10))
        
        # SELEZIONE TEMPLATE
        template_frame = ctk.CTkFrame(control_content, fg_color=("gray90", "gray20"), corner_radius=8)
        template_frame.pack(fill="x", pady=(0, 10))
        
        ctk.CTkLabel(template_frame, text="📄 Template PowerPoint:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10, 5))
        
        self.template_var = ctk.StringVar(value="Nessun template")
        self.template_combo = ctk.CTkComboBox(
            template_frame,
            variable=self.template_var,
            values=["Caricamento..."],
            command=self.on_template_selected,
            font=ctk.CTkFont(size=11),
            dropdown_font=ctk.CTkFont(size=10),
            height=32
        )
        self.template_combo.pack(fill="x", padx=10, pady=(0, 8))
        
        # Label info template selezionato
        self.template_info_label = ctk.CTkLabel(
            template_frame,
            text="ℹ️  Seleziona un template dalla lista",
            font=ctk.CTkFont(size=9),
            text_color="gray",
            justify="left"
        )
        self.template_info_label.pack(anchor="w", padx=10, pady=(0, 10))


        # SELEZIONE MODELLO LLM (semplificato)
        llm_frame = ctk.CTkFrame(control_content, fg_color=("gray90", "gray20"), corner_radius=8)
        llm_frame.pack(fill="x", pady=(0, 10))

        ctk.CTkLabel(llm_frame, text="🤖 Modello LLM:",
                    font=ctk.CTkFont(size=12, weight="bold")).pack(anchor="w", padx=10, pady=(10, 5))

        self.llm_var = ctk.StringVar(value="llama3.2:1b")
        self.llm_combo = ctk.CTkComboBox(
            llm_frame,
            variable=self.llm_var,
            values=[
                "llama3.2:1b",
                "llama3.2:3b",
            ],
            command=self.on_llm_selected,
            font=ctk.CTkFont(size=11),
            dropdown_font=ctk.CTkFont(size=10),
            height=32
        )
        self.llm_combo.pack(fill="x", padx=10, pady=(0, 10))

        # Label info modello selezionato
        #self.llm_info_label = ctk.CTkLabel(
        #    llm_frame,
        #    text="ℹ️  Modello leggero per sistemi con 8GB RAM",
        #    font=ctk.CTkFont(size=9),
        #    text_color="gray",
        #    justify="left"
        #)
        #self.llm_info_label.pack(anchor="w", padx=10, pady=(0, 10))

        # Numero candidati con slider
        candidates_frame = ctk.CTkFrame(control_content, fg_color=("gray90", "gray20"), 
                                    corner_radius=8)
        candidates_frame.pack(fill="x", pady=(0, 10))
        
        ctk.CTkLabel(candidates_frame, text="Numero candidati da estrarre:", 
                    font=ctk.CTkFont(size=12, weight="bold")).pack(pady=(10, 5))
        
        slider_container = ctk.CTkFrame(candidates_frame, fg_color="transparent")
        slider_container.pack(pady=(5, 10), padx=10)
        
        self.num_candidates = ctk.IntVar(value=1)
        self.candidates_slider = ctk.CTkSlider(slider_container, from_=1, to=10, 
                                            number_of_steps=9,
                                            variable=self.num_candidates,
                                            width=200)
        self.candidates_slider.pack(side="left", padx=(0, 10))
        
        self.num_label = ctk.CTkLabel(slider_container, text="1",
                                    font=ctk.CTkFont(size=18, weight="bold"),
                                    width=35,
                                    fg_color=("gray80", "gray25"),
                                    corner_radius=8)
        self.num_label.pack(side="left")
        self.candidates_slider.configure(command=self.update_num_label)
        
                
        # Info box
        info_frame = ctk.CTkFrame(left_column, fg_color=("gray85", "gray17"), 
                                corner_radius=8)
        info_frame.pack(fill="x", padx=10, pady=(0, 10))
        
        info_label = ctk.CTkLabel(info_frame, 
                                text="ℹ️  La ricerca utilizza AI per trovare\ni candidati più adatti alla tua gara",
                                font=ctk.CTkFont(size=10),
                                text_color="gray",
                                justify="left")
        info_label.pack(pady=8, padx=10)
        
        # ===== COLONNA DESTRA: Risultati =====
        right_column = ctk.CTkFrame(content_frame, corner_radius=10)
        right_column.grid(row=0, column=1, sticky="nsew", padx=(5, 0))
        
        # Results Header
        results_header = ctk.CTkFrame(right_column, fg_color=("gray85", "gray17"), 
                                    corner_radius=8)
        results_header.pack(fill="x", padx=10, pady=(10, 5))
        
        results_label = ctk.CTkLabel(results_header, text="📊 Risultati Pipeline",
                                    font=ctk.CTkFont(size=14, weight="bold"))
        results_label.pack(pady=8)
        
        # Results Textbox - Ridotto per fare spazio al pulsante
        self.results_text = ctk.CTkTextbox(right_column,
                                        font=ctk.CTkFont(family="Courier", size=10),
                                        corner_radius=8,
                                        wrap="word")
        self.results_text.pack(fill="both", expand=True, padx=10, pady=(0, 10))  # ← expand=True OK
        self.results_text.configure(state="disabled")

        # Search Button nella colonna destra
        button_frame = ctk.CTkFrame(right_column, fg_color="transparent")
        button_frame.pack(fill="x", padx=10, pady=(0, 10))

        self.search_button = ctk.CTkButton(button_frame, 
                                        text="🔎 Avvia Ricerca e Genera CV",
                                        command=self.run_full_pipeline,
                                        font=ctk.CTkFont(size=14, weight="bold"),
                                        height=50, 
                                        corner_radius=10,
                                        fg_color=("#3B8ED0", "#1F6AA5"),
                                        hover_color=("#36719F", "#144870"))
        self.search_button.pack(fill="x")

        # Footer info nella colonna destra
        footer_frame = ctk.CTkFrame(right_column, fg_color=("gray85", "gray17"), 
                                corner_radius=8)
        footer_frame.pack(fill="x", padx=10, pady=(0, 10))

        footer_label = ctk.CTkLabel(footer_frame, 
                                    text="💾 I risultati vengono salvati in cv_search_log.txt",
                                    font=ctk.CTkFont(size=9),
                                    text_color="gray")
        footer_label.pack(pady=6)  
    
    def parse_query_to_json(self, query_text):
        """
        Converte la query libera in un dizionario con la stessa struttura
        dei JSON dei CV. Riconosce tag come 'Skills:', 'Industry:', ecc.
        """
        query_json = {
            "name": "",
            "title": "",
            "office": "",
            "level": "",
            "summary": "",
            "skills": [],
            "technologies": [],
            "education": {"degree": "", "year": None, "program": ""},
            "certifications": [],
            "experience": []
        }

        # ── Mapping tag → campo ──────────────────────────────────
        field_map = {
            "skills":          "skills",
            "competenze":      "skills",
            "technologies":    "technologies",
            "tecnologie":      "technologies",
            "tech":            "technologies",
            "industry":        "summary",
            "settore":         "summary",
            "office":          "office",
            "sede":            "office",
            "level":           "level",
            "livello":         "level",
            "seniority":       "level",
            "role":            "title",
            "ruolo":           "title",
            "title":           "title",
            "titolo":          "title",
            "certifications":  "certifications",
            "certificazioni":  "certifications",
            "education":       "education_degree",
            "formazione":      "education_degree",
            "experience":      "experience_desc",
            "esperienza":      "experience_desc",
        }

        lines = query_text.strip().split('\n')
        unmatched_parts = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            matched = False
            for tag, field in field_map.items():
                pattern = rf'(?i)^{re.escape(tag)}\s*[:=]\s*(.+)'
                match = re.match(pattern, line)
                if match:
                    value = match.group(1).strip()
                    matched = True

                    if field in ("skills", "technologies", "certifications"):
                        items = re.split(r'[,;|]', value)
                        items = [i.strip() for i in items if i.strip()]
                        query_json[field].extend(items)

                    elif field == "summary":
                        if query_json["summary"]:
                            query_json["summary"] += " | " + value
                        else:
                            query_json["summary"] = value

                    elif field == "education_degree":
                        query_json["education"]["degree"] = value

                    elif field == "experience_desc":
                        query_json["experience"].append({
                            "company": "",
                            "period": "",
                            "description": value
                        })

                    else:
                        query_json[field] = value

                    break

            if not matched:
                unmatched_parts.append(line)

        # Testo non riconosciuto → aggiunto al summary
        if unmatched_parts:
            extra = " ".join(unmatched_parts)
            if query_json["summary"]:
                query_json["summary"] += " | " + extra
            else:
                query_json["summary"] = extra

        # Deduplica
        query_json["skills"] = list(dict.fromkeys(query_json["skills"]))
        query_json["technologies"] = list(dict.fromkeys(query_json["technologies"]))
        query_json["certifications"] = list(dict.fromkeys(query_json["certifications"]))

        return query_json

    def query_json_to_sections(self, query_json):
        """
        Converte il JSON della query nelle 4 sezioni pesate.
        
        *** IDENTICO a json_to_sections() di create_embeddings_weighted.py ***
        Questo garantisce che query e CV vivano nello stesso spazio embedding.
        """
        sections = {}

        # ── SEZIONE 1: SKILLS + TECHNOLOGIES (peso 40%) ──────────
        skills_parts = []

        skills = query_json.get("skills", [])
        if skills:
            skills_parts.append(f"Competenze tecniche: {', '.join(skills)}")

        technologies = query_json.get("technologies", [])
        if technologies:
            skills_parts.append(f"Tecnologie: {', '.join(technologies)}")

        sections['skills'] = (
            ". ".join(skills_parts) if skills_parts
            else "Nessuna competenza specificata"
        )

        # ── SEZIONE 2: EXPERIENCE (peso 40%) ─────────────────────
        experience_parts = []

        for exp in query_json.get("experience", []):
            company = exp.get("company", "")
            period = exp.get("period", "")
            description = exp.get("description", "")

            exp_text = f"Esperienza presso {company}" if company else "Esperienza"
            if period:
                exp_text += f" ({period})"
            if description:
                exp_text += f": {description}"

            experience_parts.append(exp_text)

        sections['experience'] = (
            ". ".join(experience_parts) if experience_parts
            else "Nessuna esperienza specificata"
        )

        # ── SEZIONE 3: EDUCATION + CERTIFICATIONS (peso 15%) ─────
        education_parts = []

        education = query_json.get("education", {})
        if education.get("degree"):
            edu_text = f"Formazione: {education['degree']}"
            if education.get("program"):
                edu_text += f" in {education['program']}"
            if education.get("year"):
                edu_text += f" ({education['year']})"
            education_parts.append(edu_text)

        certifications = query_json.get("certifications", [])
        if certifications:
            education_parts.append(f"Certificazioni: {', '.join(certifications)}")

        sections['education'] = (
            ". ".join(education_parts) if education_parts
            else "Nessuna formazione specificata"
        )

        # ── SEZIONE 4: SUMMARY + TITLE (peso 5%) ─────────────────
        summary_parts = []

        name = query_json.get("name", "")
        if name:
            summary_parts.append(f"Nome: {name}")

        title = query_json.get("title", "")
        if title:
            summary_parts.append(f"Ruolo: {title}")

        summary = query_json.get("summary", "")
        if summary:
            summary_parts.append(f"Profilo: {summary[:150]}")

        sections['summary'] = (
            ". ".join(summary_parts) if summary_parts
            else "Nessun sommario"
        )

        return sections

    def build_query_embedding(self, query_text):
        """
        Pipeline completa: query → JSON → 4 sezioni → 4 embeddings → media pesata.
        
        Usa gli STESSI pesi di create_embeddings_weighted.py:
          skills=40%, experience=40%, education=15%, summary=5%
        
        Ritorna (embedding, query_json, sections_dict)
        """
        # Stessi pesi di create_embeddings_weighted.py
        weights = {
            'skills': 0.40,
            'experience': 0.40,
            'education': 0.15,
            'summary': 0.05
        }

        # 1. Parse query → JSON strutturato
        query_json = self.parse_query_to_json(query_text)
        self.logger.log(f"Query JSON: {json.dumps(query_json, ensure_ascii=False)[:500]}")

        # 2. JSON → 4 sezioni (stessa funzione di create_embeddings_weighted)
        sections = self.query_json_to_sections(query_json)

        # 3. Log sezioni per debug
        self.logger.log("Query sezioni pesate:")
        for section_name, text in sections.items():
            weight_pct = weights[section_name] * 100
            self.logger.log(f"  [{weight_pct:.0f}%] {section_name}: {text[:120]}...")

        # 4. Embedding di ogni sezione separatamente
        section_embeddings = {}
        for section_name in ['skills', 'experience', 'education', 'summary']:
            emb = self.model.encode([sections[section_name]])['dense_vecs']
            section_embeddings[section_name] = emb[0]  # shape: (dim,)

        # 5. Combinazione pesata (identica a create_weighted_embeddings)
        embedding_dim = section_embeddings['skills'].shape[0]
        query_embedding = np.zeros(embedding_dim)

        for section_name, weight in weights.items():
            query_embedding += section_embeddings[section_name] * weight

        # Reshape per compatibilità con cosine_similarity: (1, dim)
        query_embedding = query_embedding.reshape(1, -1)

        self.logger.log(f"Query embedding shape: {query_embedding.shape}")

        return query_embedding, query_json, sections

      
    def load_templates(self):
        """Carica i template disponibili dalla cartella"""
        self.available_templates = get_available_templates("./template")
        
        if self.available_templates:
            template_names = [t["display_name"] for t in self.available_templates]
            self.template_combo.configure(values=template_names)
            self.template_combo.set(template_names[0])
            self.selected_template = self.available_templates[0]
            self.update_template_info()
            self.logger.log(f"Trovati {len(self.available_templates)} template")
        else:
            self.template_combo.configure(values=["Nessun template trovato"])
            self.template_combo.set("Nessun template trovato")
            self.template_info_label.configure(
                text="⚠️  Nessun template trovato in ./template/",
                text_color="orange"
            )
            self.logger.log("Nessun template trovato nella cartella ./template", "WARNING")
    
    def on_template_selected(self, choice):
        """Callback quando l'utente seleziona un template"""
        for template in self.available_templates:
            if template["display_name"] == choice:
                self.selected_template = template
                self.update_template_info()
                self.logger.log(f"Template selezionato: {template['name']}")
                break
    
    def on_llm_selected(self, choice):
        """Callback selezione modello LLM"""
        self.selected_llm_model = choice.strip()
        self.logger.log(f"Modello LLM selezionato: {self.selected_llm_model}")

        
    def update_template_info(self):
        """Aggiorna le informazioni del template selezionato"""
        if self.selected_template:
            info_text = f"✓ Template: {self.selected_template['name']}\n"
            info_text += f"📁 Percorso: {self.selected_template['path'].name}"
            self.template_info_label.configure(
                text=info_text,
                text_color=("#2CC985" if ctk.get_appearance_mode() == "Dark" else "#1B7F5A")
            )
    
    def update_num_label(self, value):
        """Aggiorna label numero candidati"""
        self.num_label.configure(text=str(int(value)))
    
    def load_data(self):
        """Carica dati NPY subito, modello BGE-M3 in background"""
        try:
            required_files = ['cv_embeddings.npy', 'cv_texts.npy', 'cv_labels.npy']
            missing = [f for f in required_files if not os.path.exists(f)]

            if missing:
                messagebox.showerror("Errore",
                    f"File mancanti: {', '.join(missing)}\n\nEsegui prima create_embeddings_weighted.py")
                self.status_label.configure(text="❌ File mancanti!", text_color="red")
                self.logger.log(f"File mancanti: {missing}", "ERROR")
                return

            # Carica subito i file NPY (veloce, ~100ms)
            self.cv_embeddings = np.load('cv_embeddings.npy')
            self.cv_texts = np.load('cv_texts.npy', allow_pickle=True)
            self.cv_labels = np.load('cv_labels.npy', allow_pickle=True)

            self.status_label.configure(
                text=f"⏳ {len(self.cv_labels)} CV caricati — modello in caricamento...",
                text_color="orange")
            self.root.update()

            # Carica il modello pesante in un thread separato
            self.search_button.configure(state="disabled")
            thread = threading.Thread(target=self._load_model_background, daemon=True)
            thread.start()

        except Exception as e:
            messagebox.showerror("Errore", str(e))
            self.status_label.configure(text="❌ Errore caricamento", text_color="red")
            self.logger.log(f"Errore caricamento: {e}", "ERROR")

    def _load_model_background(self):
        """Carica BGE-M3 in background senza bloccare la UI"""
        try:
            self.model = BGEM3FlagModel('BAAI/bge-m3', use_fp16=True)
            self.logger.log("Modello BGE-M3 caricato (background)")

            # Aggiorna UI dal thread principale
            self.root.after(0, self._on_model_ready)

        except Exception as e:
            self.logger.log(f"Errore caricamento modello: {e}", "ERROR")
            self.root.after(0, lambda: self.status_label.configure(
                text="❌ Errore caricamento modello", text_color="red"))

    def _on_model_ready(self):
        """Callback quando il modello è pronto"""
        self.status_label.configure(
            text=f"✅ Sistema pronto! {len(self.cv_labels)} CV caricati",
            text_color="#2CC985")
        self.search_button.configure(state="normal")
        self.logger.log(f"Sistema pronto: {len(self.cv_labels)} CV caricati")

    
    def find_existing_json(self, label, json_folder):
        """Cerca un JSON esistente"""
        json_folder = Path(json_folder)
        if not json_folder.exists():
            return None
        
        normalized = label.lower().strip().replace(' ', '_').replace('.', '_')
        
        for json_file in json_folder.glob("*.json"):
            filename_lower = json_file.stem.lower()
            if normalized == filename_lower:
                return json_file
            
            name_parts = normalized.split('_')
            if all(part in filename_lower for part in name_parts if part):
                return json_file
        
        return None
    
    def analyze_cv_with_llm(self, cv_data, query, similarity_score):
        """Analizza CV usando Ollama LLM locale"""
        try:
            prompt = f"""Analizza questo CV per la gara.

GARA: {query[:200]}

CV:
- Nome: {cv_data.get('name', 'N/A')}
- Ruolo: {cv_data.get('title', 'N/A')}
- Skills: {', '.join(cv_data.get('skills', [])[:8])}
- Tech: {', '.join(cv_data.get('technologies', [])[:8])}
- Score: {similarity_score:.2f}

Fornisci (max 200 parole):
1. VALUTAZIONE: <punteggio 0-100>
2. PUNTI FORZA: <2-3 bullet points>
3. GAP: <eventuali lacune>
4. ESITO: IDONEO/DA_VALUTARE/NON_IDONEO

Rispondi in italiano, formato chiaro."""

            # Timeout aumentato per primo caricamento modello
            response = requests.post(
                'http://localhost:11434/api/generate',
                json={
                    "model": self.selected_llm_model,
                    "prompt": prompt,
                    "stream": False,
                    "options": {
                    "temperature": 0.7,
                    "num_predict": 512,      # ← AUMENTATO (era 200)
                    "num_ctx": 2048,         # ← Contesto più grande
                    "stop": ["---", "###"]   # ← Stop solo su delimitatori specifici
                }
                },
                timeout=180  # ← AUMENTATO A 3 MINUTI
            )
            
            if response.status_code == 200:
                return response.json()['response']
            else:
                return f"⚠️ Errore LLM: status {response.status_code}"
                
        except requests.exceptions.ConnectionError:
            return "⚠️ Ollama non disponibile. Verifica che sia in esecuzione."
        except requests.exceptions.Timeout:
            return "⚠️ Timeout LLM: il modello sta caricando, riprova tra 1 minuto."
        except Exception as e:
            return f"⚠️ Errore analisi LLM: {str(e)[:100]}"

    def plot_pca_3d(self, query_embedding, top_indices, similarities):
        """Visualizza grafico 3D PCA con query e candidati"""
        try:
            # Combina embeddings: query + tutti i CV
            all_embeddings = np.vstack([query_embedding, self.cv_embeddings])
            
            # Applica PCA per ridurre a 3 dimensioni
            pca = PCA(n_components=3)
            embeddings_3d = pca.fit_transform(all_embeddings)
            
            # Separa query (primo punto) dagli altri
            query_3d = embeddings_3d[0]
            cv_embeddings_3d = embeddings_3d[1:]
            
            # Crea figura
            fig = plt.figure(figsize=(12, 9))
            ax = fig.add_subplot(111, projection='3d')
            
            # Plot tutti i CV (grigi e piccoli)
            ax.scatter(cv_embeddings_3d[:, 0], 
                      cv_embeddings_3d[:, 1], 
                      cv_embeddings_3d[:, 2],
                      c='lightgray', 
                      marker='o', 
                      s=20, 
                      alpha=0.3,
                      edgecolors='black',    # ← Bordo grigio scuro
                      linewidth=2.0,         # ← Spessore bordo
                      label='Altri CV')
            
            # Plot top candidati (colorati per similarità)
            top_embeddings_3d = cv_embeddings_3d[top_indices]
            top_similarities = similarities[top_indices]
            
            scatter = ax.scatter(top_embeddings_3d[:, 0],
                                top_embeddings_3d[:, 1],
                                top_embeddings_3d[:, 2],
                                c=top_similarities,
                                cmap='RdYlGn',
                                marker='o',
                                s=150,
                                alpha=0.8,
                                edgecolors='black',
                                linewidth=1.5,
                                label='Top Candidati')
            
            # Aggiungi etichette per i top candidati
            for idx, (i, sim) in enumerate(zip(top_indices, top_similarities)):
                label = self.cv_labels[i]
                ax.text(top_embeddings_3d[idx, 0],
                       top_embeddings_3d[idx, 1],
                       top_embeddings_3d[idx, 2],
                       f'{label}\n({sim:.3f})',
                       fontsize=8,
                       bbox=dict(boxstyle='round,pad=0.3', facecolor='yellow', alpha=0.7))
            
            # Plot query (stella rossa grande)
            ax.scatter(query_3d[0], query_3d[1], query_3d[2],
                      c='red',
                      marker='*',
                      s=500,
                      edgecolors='black',
                      linewidth=2,
                      label='Query',
                      zorder=1000)
            
            ax.text(query_3d[0], query_3d[1], query_3d[2],
                   '  QUERY',
                   fontsize=12,
                   fontweight='bold',
                   color='red')
            
            # Configura assi
            ax.set_xlabel(f'PC1 ({pca.explained_variance_ratio_[0]*100:.1f}%)', fontsize=10)
            ax.set_ylabel(f'PC2 ({pca.explained_variance_ratio_[1]*100:.1f}%)', fontsize=10)
            ax.set_zlabel(f'PC3 ({pca.explained_variance_ratio_[2]*100:.1f}%)', fontsize=10)
            ax.set_title('Visualizzazione 3D PCA: Query e Candidati\nVarianza spiegata: {:.1f}%'.format(
                sum(pca.explained_variance_ratio_) * 100), fontsize=14, fontweight='bold')
            
            # Colorbar
            cbar = plt.colorbar(scatter, ax=ax, pad=0.1, shrink=0.8)
            cbar.set_label('Similarità con Query', fontsize=10)
            
            # Leggenda
            ax.legend(loc='upper left', fontsize=10)
            
            # Griglia
            ax.grid(True, alpha=0.3)
            
            plt.tight_layout()
            plt.show()
            
            self.logger.log("Grafico PCA 3D generato con successo")
            
        except Exception as e:
            self.logger.log(f"Errore generazione grafico PCA: {e}", "ERROR")
            messagebox.showerror("Errore", f"Impossibile generare il grafico PCA:\n{e}")


    def generate_direct_cv(self):
        """Genera CV direttamente da nome e cognome senza ricerca"""
        name_input = self.direct_name_entry.get().strip()
        
        if not name_input:
            messagebox.showwarning("Attenzione", "Inserisci nome e cognome!")
            return
        
        if not self.selected_template:
            messagebox.showerror("Errore", "Seleziona un template!")
            return
        
        try:
            # Disabilita bottone durante elaborazione
            self.direct_gen_button.configure(state="disabled", text="⏳ Generazione...")
            self.root.update()
            
            self.logger.log(f"=== GENERAZIONE DIRETTA CV ===")
            self.logger.log(f"Nome: {name_input}")
            self.logger.log(f"Template: {self.selected_template['name']}")
            
            self.append_result(f"\n{'═'*80}\n")
            self.append_result(f"⚡ GENERAZIONE DIRETTA CV\n")
            self.append_result(f"{'═'*80}\n")
            self.append_result(f"Candidato: {name_input}\n")
            self.append_result(f"Template: {self.selected_template['name']}\n\n")
            
            # STEP 1: Cerca JSON esistente
            self.status_label.configure(text="⏳ Ricerca JSON...", text_color="orange")
            self.root.update()
            
            self.append_result("📁 STEP 1: RICERCA JSON\n")
            self.append_result("─"*80 + "\n")
            
            extractor = PPTXToJSONExtractor(logger=self.logger)
            existing_json = self.find_existing_json(name_input, extractor.cv_json_folder)
            
            if existing_json and existing_json.exists():
                json_file = existing_json
                self.append_result(f"✅ JSON trovato: {existing_json.name}\n\n")
            else:
                # Prova estrazione da PPTX
                self.append_result(f"→ JSON non trovato, estrazione da PPTX...\n")
                json_file = extractor.process_label(name_input)
                
                if not json_file or not json_file.exists():
                    self.append_result(f"❌ ERRORE: Nessun PPTX trovato per '{name_input}'\n")
                    messagebox.showerror("Errore", 
                        f"Non è stato trovato né un JSON né un PPTX per '{name_input}'.\n\n"
                        f"Verifica che il file esista in:\n"
                        f"- {extractor.cv_ppt_folder}\n"
                        f"- {extractor.cv_json_folder}")
                    self.direct_gen_button.configure(state="normal", text="⚡ Genera CV Diretto")
                    self.status_label.configure(text="❌ File non trovato", text_color="red")
                    return
                
                self.append_result(f"✅ JSON estratto: {json_file.name}\n\n")
            
            # STEP 2: Generazione CV
            self.status_label.configure(text="⏳ Generazione CV...", text_color="orange")
            self.root.update()
            
            self.append_result("📄 STEP 2: GENERAZIONE CV\n")
            self.append_result("─"*80 + "\n")
            
            output_folder = Path("./cv_generati")
            output_folder.mkdir(exist_ok=True)
            
            json_stem = json_file.stem
            if json_stem.lower().startswith("cv_"):
                json_stem = json_stem[3:]
            
            output_file = output_folder / f"CV_{json_stem}.pptx"
            
            # Crea generator
            generator = create_generator_for_template(self.selected_template['path'], self.logger)
            
            if generator.generate_cv(json_file, output_file):
                self.append_result(f"✅ CV generato: {output_file.name}\n\n")
                
                # RIEPILOGO
                self.append_result("\n" + "═"*80 + "\n")
                self.append_result("✨ GENERAZIONE COMPLETATA\n")
                self.append_result("═"*80 + "\n")
                self.append_result(f"✅ File generato: {output_file.name}\n")
                self.append_result(f"📁 Cartella: {output_folder}\n")
                self.append_result("═"*80 + "\n")
                
                self.status_label.configure(text=f"✅ CV generato con successo!", 
                                           text_color="#2CC985")
                self.logger.log(f"=== CV GENERATO: {output_file.name} ===")
                
                messagebox.showinfo("✅ Successo!", 
                    f"CV generato con successo!\n\n"
                    f"File: {output_file.name}\n"
                    f"Cartella: {output_folder}")
            else:
                self.append_result(f"❌ Errore nella generazione del CV\n")
                messagebox.showerror("Errore", "Errore durante la generazione del CV")
                self.status_label.configure(text="❌ Errore generazione", text_color="red")
            
            # Riabilita bottone
            self.direct_gen_button.configure(state="normal", text="⚡ Genera CV Diretto")
            
        except Exception as e:
            error_msg = f"Errore generazione diretta: {e}"
            self.append_result(f"\n❌ ERRORE: {error_msg}\n")
            messagebox.showerror("Errore", error_msg)
            self.status_label.configure(text="❌ Errore generazione", text_color="red")
            self.logger.log(error_msg, "ERROR")
            self.direct_gen_button.configure(state="normal", text="⚡ Genera CV Diretto")

    def run_full_pipeline(self):
        """Esegue la pipeline completa"""
        if self.model is None:
            messagebox.showwarning("Attenzione", "Modello non caricato")
            return
        
        query = self.query_text.get("1.0", "end").strip()
        if not query:
            messagebox.showwarning("Attenzione", "Inserisci una query!")
            return
        
        try:
            # Disabilita bottone durante elaborazione
            self.search_button.configure(state="disabled", text="⏳ Elaborazione...")
            self.root.update()
            
            num_candidates = int(self.num_candidates.get())
            self.logger.log(f"=== INIZIO PIPELINE ===")
            self.logger.log(f"Query: {query[:100]}...")
            self.logger.log(f"Numero candidati: {num_candidates}")
            self.logger.log(f"Modello LLM: {self.selected_llm_model}")
            
            self.append_result(f"\n{'═'*80}\n")
            self.append_result(f"🚀 AVVIO PIPELINE COMPLETA\n")
            self.append_result(f"{'═'*80}\n")
            self.append_result(f"Candidati richiesti: {num_candidates}\n\n")
            
            # STEP 1
            self.status_label.configure(text="⏳ Step 1/3: Calcolo similarità...", text_color="orange")
            self.root.update()
            
            # Usa lo stesso processo pesato di create_embeddings_weighted.py
            query_embedding, query_json, query_sections = self.build_query_embedding(query)
            similarities = cosine_similarity(query_embedding, self.cv_embeddings)[0]
            
            # Mostra nei risultati le sezioni pesate
            self.append_result("🔄 QUERY NORMALIZZATA (Weighted Sections):\n")
            self.append_result("─"*80 + "\n")
            weights_display = {'skills': '40%', 'experience': '40%', 'education': '15%', 'summary': '5%'}
            for section, text in query_sections.items():
                self.append_result(f"  [{weights_display[section]}] {section}: {text[:100]}\n")
                self.append_result("\n")
            sorted_indices = np.argsort(similarities)[::-1]
            top_candidates = sorted_indices[:num_candidates]
            
            self.append_result("📊 STEP 1: CANDIDATI SELEZIONATI\n")
            self.append_result("─"*80 + "\n")
            
            selected_labels = []
            for rank, idx in enumerate(top_candidates, 1):
                label = self.cv_labels[idx]
                sim = similarities[idx]
                selected_labels.append(label)
                self.append_result(f"  {rank}. {label}\n     Similarità: {sim:.4f} ({sim*100:.2f}%)\n\n")
                self.logger.log(f"Candidato {rank}: {label} (sim={sim:.4f})")

            # ANALISI LLM PER OGNI CANDIDATO
            self.append_result("\n🤖 ANALISI LLM DEI CANDIDATI\n")
            self.append_result("─"*80 + "\n")
            self.status_label.configure(text="⏳ Analisi LLM in corso...", text_color="orange")

            extractor_temp = PPTXToJSONExtractor(logger=self.logger)

            for rank, (idx, label) in enumerate(zip(top_candidates, selected_labels), 1):
                # Carica JSON del candidato
                json_file = self.find_existing_json(label, extractor_temp.cv_json_folder)
                
                if json_file and json_file.exists():
                    with open(json_file, 'r', encoding='utf-8') as f:
                        cv_data = json.load(f)
                    
                    # Analisi LLM
                    self.append_result(f"\n[{rank}] {label}:\n")
                    self.root.update()
                    
                    analysis = self.analyze_cv_with_llm(cv_data, query, similarities[idx])
                    self.append_result(f"{analysis}\n")
                    self.append_result("─"*40 + "\n")
                    self.logger.log(f"Analisi LLM completata per: {label}")
                else:
                    self.append_result(f"\n[{rank}] {label}: ⚠️ JSON non disponibile per analisi\n")

            # Visualizza grafico 3D PCA (FUORI DAL LOOP)
            self.append_result("\n📈 Generazione grafico PCA 3D...\n")
            self.root.update()
            self.plot_pca_3d(query_embedding, top_candidates, similarities)
            
            # STEP 2
            self.status_label.configure(text="⏳ Step 2/3: Verifica JSON...", text_color="orange")
            self.root.update()
            
            self.append_result("\n📁 STEP 2: VERIFICA JSON\n")
            self.append_result("─"*80 + "\n")
            
            extractor = PPTXToJSONExtractor(logger=self.logger)
            json_files = []
            
            for i, label in enumerate(selected_labels, 1):
                self.append_result(f"[{i}/{len(selected_labels)}] {label}...\n")
                self.root.update()
                
                existing_json = self.find_existing_json(label, extractor.cv_json_folder)
                
                if existing_json and existing_json.exists():
                    json_files.append(existing_json)
                    self.append_result(f"  ✅ JSON esistente: {existing_json.name}\n\n")
                else:
                    self.append_result(f"  → Estrazione da PPTX...\n")
                    json_file = extractor.process_label(label)
                    
                    if json_file and json_file.exists():
                        json_files.append(json_file)
                        self.append_result(f"  ✅ JSON creato: {json_file.name}\n\n")
                    else:
                        self.append_result(f"  ❌ ERRORE: PPTX non trovato\n\n")
            
            # STEP 3
            self.status_label.configure(text="⏳ Step 3/3: Generazione CV...", text_color="orange")
            self.root.update()
            
            self.append_result("\n📄 STEP 3: GENERAZIONE CV\n")
            self.append_result("─"*80 + "\n")
            
            # Verifica che sia stato selezionato un template
            if not self.selected_template:
                messagebox.showerror("Errore", "Nessun template selezionato!")
                self.search_button.configure(state="normal", text="🔎 Avvia Ricerca e Genera CV")
                return

            output_folder = Path("./cv_generati")
            output_folder.mkdir(exist_ok=True)
            
            generated_files = []
            for i, json_file in enumerate(json_files, 1):
                json_stem = json_file.stem
                if json_stem.lower().startswith("cv_"):
                    json_stem = json_stem[3:]
                
                output_file = output_folder / f"CV_{json_stem}.pptx"
                self.append_result(f"[{i}/{len(json_files)}] {json_stem}...\n")
                
                # Crea il generator per il template selezionato
                generator = create_generator_for_template(self.selected_template['path'], self.logger)

                if generator.generate_cv(json_file, output_file):
                    generated_files.append(output_file)
                    self.append_result(f"  ✅ CV generato: {output_file.name}\n\n")
                else:
                    self.append_result(f"  ❌ Errore nella generazione\n\n")
            
            # RIEPILOGO
            self.append_result("\n" + "═"*80 + "\n")
            self.append_result("✨ PIPELINE COMPLETATA\n")
            self.append_result("═"*80 + "\n")
            self.append_result(f"✅ Candidati processati: {len(selected_labels)}\n")
            self.append_result(f"✅ JSON generati/riutilizzati: {len(json_files)}\n")
            self.append_result(f"✅ CV generati: {len(generated_files)}\n")
            self.append_result(f"📁 Cartella output: {output_folder}\n")
            self.append_result("═"*80 + "\n")
            
            self.status_label.configure(text=f"✅ Pipeline completata! {len(generated_files)} CV generati", 
                                       text_color="#2CC985")
            self.logger.log(f"=== PIPELINE COMPLETATA: {len(generated_files)} CV generati ===")
            
            # Riabilita bottone
            self.search_button.configure(state="normal", text="🔎 Cerca e Genera CV")
            
            messagebox.showinfo("✅ Successo!", 
                f"Pipeline completata con successo!\n\n"
                f"CV generati: {len(generated_files)}\n"
                f"Cartella: {output_folder}\n\n"
                f"Log dettagliato: cv_search_log.txt")
            
        except Exception as e:
            error_msg = f"Errore nella pipeline: {e}"
            self.append_result(f"\n❌ ERRORE: {error_msg}\n")
            messagebox.showerror("Errore", error_msg)
            self.status_label.configure(text="❌ Errore nella pipeline", text_color="red")
            self.logger.log(error_msg, "ERROR")
            self.search_button.configure(state="normal", text="🔎 Cerca e Genera CV")
    
    def append_result(self, text):
        """Aggiunge testo ai risultati"""
        self.results_text.configure(state="normal")
        self.results_text.insert("end", text)
        self.results_text.see("end")
        self.results_text.configure(state="disabled")
        self.root.update()
    
    def run(self):
        """Avvia l'applicazione"""
        self.root.mainloop()

def main():
    app = CVSearchApp()
    app.run()

if __name__ == "__main__":
    main()                    