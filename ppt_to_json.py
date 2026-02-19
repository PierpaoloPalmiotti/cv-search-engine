"""
Estrattore di informazioni da PowerPoint a JSON strutturato
Step 2 del processo: da nome persona -> PPTX in cv_ppt -> JSON in cv_json
"""

import json
import re
from pathlib import Path
from typing import List, Dict, Optional
from pptx import Presentation
from PIL import Image
import pytesseract
import io

class PPTXToJSONExtractor:
    def __init__(self, cv_ppt_folder: str = "./cv_ppt", 
                 cv_json_folder: str = "./cv_json",
                 use_ocr: bool = False, 
                 ocr_lang: str = 'ita+eng'):
        """
        Args:
            cv_ppt_folder: Cartella contenente i file PPTX
            cv_json_folder: Cartella dove salvare i JSON
            use_ocr: Se True, usa OCR anche per immagini nelle slide
            ocr_lang: Lingue per OCR (default: italiano + inglese)
        """
        self.cv_ppt_folder = Path(cv_ppt_folder)
        self.cv_json_folder = Path(cv_json_folder)
        self.use_ocr = use_ocr
        self.ocr_lang = ocr_lang
        
        # Crea cartella output se non esiste
        self.cv_json_folder.mkdir(exist_ok=True, parents=True)
        
    def extract_text_from_shapes(self, slide) -> str:
        """Estrae tutto il testo dalle shape della slide"""
        all_text = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                all_text.append(shape.text.strip())
            
            if shape.shape_type == 19:  # Table
                table_text = self._extract_from_table(shape)
                if table_text:
                    all_text.append(table_text)
            
            if self.use_ocr and self._is_image_shape(shape):
                ocr_text = self._extract_text_from_image_shape(shape)
                if ocr_text:
                    all_text.append(ocr_text)
        
        return "\n".join(all_text)
    
    def _extract_from_table(self, shape) -> str:
        """Estrae testo da tabelle"""
        table_content = []
        try:
            table = shape.table
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.strip():
                    table_content.append(row_text)
        except:
            pass
        return "\n".join(table_content)
    
    def _is_image_shape(self, shape) -> bool:
        """Verifica se la shape contiene un'immagine"""
        return hasattr(shape, "image") or shape.shape_type == 13
    
    def _extract_text_from_image_shape(self, shape) -> str:
        """Estrae testo da immagine usando OCR"""
        try:
            if hasattr(shape, "image"):
                image_bytes = shape.image.blob
                image = Image.open(io.BytesIO(image_bytes))
                text = pytesseract.image_to_string(image, lang=self.ocr_lang)
                return text.strip()
        except Exception as e:
            print(f"  [OCR Warning] {e}")
        return ""
    
    def find_pptx_by_name(self, person_name: str) -> Optional[Path]:
        """
        Trova il file PPTX corrispondente al nome della persona
        Supporta vari formati: nome_cognome.pptx, nome.cognome.pptx, ecc.
        """
        if not self.cv_ppt_folder.exists():
            print(f"[Error] Cartella cv_ppt non trovata: {self.cv_ppt_folder}")
            return None
        
        # Normalizza il nome per la ricerca
        normalized_name = person_name.lower().strip()
        normalized_name = re.sub(r'\s+', '_', normalized_name)  # Spazi -> underscore
        
        # Pattern di ricerca possibili
        search_patterns = [
            f"{normalized_name}.pptx",
            f"{normalized_name.replace('_', '.')}.pptx",
            f"{normalized_name.replace('_', '-')}.pptx",
            f"{normalized_name.replace('_', '')}.pptx",
        ]
        
        # Cerca file con match esatto
        for pattern in search_patterns:
            potential_file = self.cv_ppt_folder / pattern
            if potential_file.exists():
                return potential_file
        
        # Cerca con match parziale (case-insensitive)
        name_parts = normalized_name.replace('_', ' ').split()
        for pptx_file in self.cv_ppt_folder.glob("*.pptx"):
            filename_lower = pptx_file.stem.lower()
            # Verifica se tutte le parti del nome sono nel filename
            if all(part in filename_lower for part in name_parts):
                return pptx_file
        
        print(f"[Warning] Nessun PPTX trovato per: {person_name}")
        return None
    
    def extract_name(self, text: str) -> str:
        """Estrae il nome completo"""
        lines = text.split('\n')
        
        patterns = [
            r'^([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)',
            r'(?:Nome|Name)[:\s]+([A-Z][a-z]+(?:\s+[A-Z][a-z]+)+)',
            r'([A-Z][A-Z\s]+)',
        ]
        
        for line in lines[:5]:
            for pattern in patterns:
                match = re.search(pattern, line)
                if match:
                    name = match.group(1).strip()
                    if len(name.split()) >= 2 and len(name) > 5:
                        if not any(kw in name.lower() for kw in ['curriculum', 'profilo', 'competenze', 'profile']):
                            return name.title()
        
        return ""
    
    def extract_title(self, text: str) -> str:
        """Estrae il titolo/ruolo professionale"""
        keywords = ['ruolo', 'posizione', 'title', 'job title', 'position', 'role']
        lines = text.split('\n')
        
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            
            if any(kw in line_lower for kw in keywords):
                title = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if title:
                    return title
                if i + 1 < len(lines):
                    return lines[i + 1].strip()
            
            if i > 0 and i < 5:
                if re.search(r'\b(senior|junior|lead|manager|developer|engineer|analyst|consultant|specialist)\b', line_lower):
                    return line.strip()
        
        return ""
    
    def extract_summary(self, text: str) -> str:
        """Estrae il sommario/profilo"""
        keywords = ['profilo', 'summary', 'about', 'bio', 'descrizione', 'profile', 'chi sono', 'presentazione']
        lines = text.split('\n')
        summary_lines = []
        capturing = False
        
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            
            if any(kw in line_lower for kw in keywords):
                capturing = True
                content = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if content and len(content) > 20:
                    summary_lines.append(content)
                continue
            
            if capturing:
                if any(stop in line_lower for stop in ['competenze', 'formazione', 'esperienza', 'certificazioni', 'skills', 'education', 'experience']):
                    break
                
                if line.strip() and len(line.strip()) > 10:
                    summary_lines.append(line.strip())
                    
                if len(summary_lines) >= 5:
                    break
        
        return " ".join(summary_lines)
    
    def extract_list_items(self, text: str, keywords: List[str]) -> List[str]:
        """Estrae lista di elementi (skills, technologies, certificazioni)"""
        lines = text.split('\n')
        items = []
        capturing = False
        
        stop_keywords = ['formazione', 'esperienza', 'certificazioni', 'competenze', 
                        'education', 'experience', 'certifications', 'skills', 'tecnologie']
        
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            
            if any(kw in line_lower for kw in keywords):
                capturing = True
                content = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if content:
                    items.extend(self._parse_items(content))
                continue
            
            if capturing:
                is_other_section = any(stop in line_lower for stop in stop_keywords)
                is_current_section = any(kw in line_lower for kw in keywords)
                
                if is_other_section and not is_current_section:
                    break
                
                if line.strip() and len(line.strip()) > 1:
                    items.extend(self._parse_items(line))
        
        seen = set()
        unique_items = []
        for item in items:
            item_lower = item.lower()
            if item_lower not in seen and len(item) > 1:
                seen.add(item_lower)
                unique_items.append(item)
        
        return unique_items
    
    def _parse_items(self, text: str) -> List[str]:
        """Parse una riga in items individuali"""
        text = re.sub(r'^[\-\•\*\d\.\)]+\s*', '', text.strip())
        
        separators = [',', ';', '|', '–', '•']
        for sep in separators:
            if sep in text:
                items = [item.strip() for item in text.split(sep)]
                return [item for item in items if len(item) > 1]
        
        if len(text) > 1:
            return [text]
        
        return []
    
    def extract_education(self, text: str) -> Dict:
        """Estrae informazioni sulla formazione"""
        keywords = ['formazione', 'istruzione', 'education', 'laurea', 'degree', 'studi']
        lines = text.split('\n')
        
        education = {
            "degree": "",
            "year": None,
            "program": ""
        }
        
        capturing = False
        edu_lines = []
        
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            
            if any(kw in line_lower for kw in keywords):
                capturing = True
                content = re.sub(r'(' + '|'.join(keywords) + r')[:\s]*', '', line, flags=re.IGNORECASE).strip()
                if content:
                    edu_lines.append(content)
                continue
            
            if capturing:
                if any(stop in line_lower for stop in ['esperienza', 'competenze', 'certificazioni', 'experience', 'skills']):
                    break
                
                if line.strip():
                    edu_lines.append(line.strip())
                    
                if len(edu_lines) >= 3:
                    break
        
        full_text = " ".join(edu_lines)
        
        year_match = re.search(r'\b(19|20)\d{2}\b', full_text)
        if year_match:
            education["year"] = int(year_match.group(0))
        
        degree_patterns = [
            r'(Laurea\s+(?:Magistrale|Triennale)?(?:\s+in\s+[^,\.\n]+)?)',
            r'(Master\s+(?:of\s+Science|in\s+[^,\.\n]+)?)',
            r'(Bachelor\s+(?:of\s+Science|in\s+[^,\.\n]+)?)',
            r'(Diploma\s+[^,\.\n]+)',
            r'(PhD|Dottorato)(?:\s+in\s+[^,\.\n]+)?'
        ]
        
        for pattern in degree_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE)
            if match:
                education["degree"] = match.group(1).strip()
                break
        
        if not education["degree"] and edu_lines:
            education["program"] = edu_lines[0]
        else:
            program_match = re.search(r'(?:in|corso)\s+([A-Z][^,\.\n]{5,})', full_text)
            if program_match:
                education["program"] = program_match.group(1).strip()
        
        return education
    
    def extract_experience(self, text: str) -> List[Dict]:
        """Estrae esperienze lavorative"""
        keywords = ['esperienza', 'experience', 'lavoro', 'work', 'carriera', 'employment']
        lines = text.split('\n')
        
        experiences = []
        capturing = False
        current_exp = []
        
        for i, line in enumerate(lines):
            line_lower = line.lower().strip()
            
            if any(kw in line_lower for kw in keywords):
                capturing = True
                continue
            
            if capturing:
                if any(stop in line_lower for stop in ['formazione', 'competenze', 'certificazioni', 'education', 'skills']):
                    break
                
                if line.strip():
                    if re.search(r'\b(19|20)\d{2}\b|(?:gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic|jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)', line_lower):
                        if current_exp:
                            experiences.append(self._parse_experience(current_exp))
                            current_exp = []
                    
                    current_exp.append(line.strip())
        
        if current_exp:
            experiences.append(self._parse_experience(current_exp))
        
        return experiences
    
    def _parse_experience(self, lines: List[str]) -> Dict:
        """Parse singola esperienza"""
        full_text = " ".join(lines)
        
        experience = {
            "company": "",
            "period": "",
            "description": ""
        }
        
        period_patterns = [
            r'(\d{4}\s*-\s*\d{4})',
            r'(\d{4}\s*-\s*(?:Presente|Present|oggi|today))',
            r'((?:gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic)\s+\d{4}\s*-\s*(?:gen|feb|mar|apr|mag|giu|lug|ago|set|ott|nov|dic)\s+\d{4})'
        ]
        
        for pattern in period_patterns:
            match = re.search(pattern, full_text, re.IGNORECASE)
            if match:
                experience["period"] = match.group(1).strip()
                break
        
        if lines:
            experience["company"] = lines[0]
        
        if len(lines) > 1:
            experience["description"] = " ".join(lines[1:])
        
        return experience
    
    def extract_info_from_pptx(self, pptx_path: Path) -> Dict:
        """Estrae tutte le informazioni da un file PPTX e crea JSON strutturato"""
        try:
            prs = Presentation(str(pptx_path))
            
            if len(prs.slides) == 0:
                print(f"  [Warning] Nessuna slide trovata in {pptx_path.name}")
                return self._empty_json()
            
            slide = prs.slides[0]
            all_text = self.extract_text_from_shapes(slide)
            
            if not all_text.strip():
                print(f"  [Warning] Nessun testo estratto da {pptx_path.name}")
                return self._empty_json()
            
            json_data = {
                "name": self.extract_name(all_text),
                "title": self.extract_title(all_text),
                "summary": self.extract_summary(all_text),
                "skills": self.extract_list_items(all_text, 
                    ['competenze', 'skills', 'abilità', 'capacità']),
                "technologies": self.extract_list_items(all_text, 
                    ['tecnologie', 'technologies', 'tech stack', 'strumenti', 'tools']),
                "education": self.extract_education(all_text),
                "certifications": self.extract_list_items(all_text, 
                    ['certificazioni', 'certifications', 'certificati']),
                "experience": self.extract_experience(all_text)
            }
            
            return json_data
            
        except Exception as e:
            print(f"  [Error] Impossibile processare {pptx_path.name}: {e}")
            return self._empty_json()
    
    def _empty_json(self) -> Dict:
        """Restituisce JSON vuoto con struttura corretta"""
        return {
            "name": "",
            "title": "",
            "summary": "",
            "skills": [],
            "technologies": [],
            "education": {
                "degree": "",
                "year": None,
                "program": ""
            },
            "certifications": [],
            "experience": []
        }
    
    def process_person(self, person_name: str) -> Optional[Dict]:
        """
        Step 2 del processo: dato il nome, trova il PPTX e crea il JSON
        
        Args:
            person_name: Nome della persona (output dello step 1)
            
        Returns:
            Dict con i dati estratti o None se fallisce
        """
        print(f"\n[Step 2] Processando: {person_name}")
        
        # 1. Trova il file PPTX
        pptx_file = self.find_pptx_by_name(person_name)
        if not pptx_file:
            return None
        
        print(f"  ✓ PPTX trovato: {pptx_file.name}")
        
        # 2. Estrai informazioni
        json_data = self.extract_info_from_pptx(pptx_file)
        
        # 3. Salva JSON
        json_filename = self.cv_json_folder / f"{pptx_file.stem}.json"
        with open(json_filename, 'w', encoding='utf-8') as f:
            json.dump(json_data, f, ensure_ascii=False, indent=2)
        
        print(f"  ✓ JSON salvato: {json_filename.name}")
        
        return json_data
    
    def process_person_list(self, person_names: List[str]) -> List[Dict]:
        """
        Processa una lista di nomi (output dello step 1)
        
        Args:
            person_names: Lista di nomi da processare
            
        Returns:
            Lista di dict con i dati estratti
        """
        print(f"\n{'='*60}")
        print(f"STEP 2: Estrazione da PPTX -> JSON")
        print(f"{'='*60}")
        print(f"Persone da processare: {len(person_names)}")
        print(f"Cartella PPTX: {self.cv_ppt_folder}")
        print(f"Cartella JSON: {self.cv_json_folder}")
        print(f"{'='*60}\n")
        
        results = []
        success_count = 0
        
        for i, person_name in enumerate(person_names, 1):
            print(f"[{i}/{len(person_names)}] {person_name}")
            
            json_data = self.process_person(person_name)
            
            if json_data:
                results.append(json_data)
                success_count += 1
        
        print(f"\n{'='*60}")
        print(f"✅ Completato! {success_count}/{len(person_names)} profili elaborati")
        print(f"📁 JSON salvati in: {self.cv_json_folder}")
        print(f"{'='*60}\n")
        
        return results


# ============ ESEMPI DI UTILIZZO ============

def main():
    """Esempi di utilizzo come Step 2 del processo"""
    
    # Inizializza l'estrattore con le cartelle corrette
    extractor = PPTXToJSONExtractor(
        cv_ppt_folder="./cv_ppt",
        cv_json_folder="./cv_json",
        use_ocr=False
    )
    
    # CASO 1: Output dello Step 1 è un singolo nome
    person_name = "Mario Rossi"  # Output dello Step 1
    extractor.process_person(person_name)
    
    # CASO 2: Output dello Step 1 è una lista di nomi
    # person_names = ["Mario Rossi", "Giulia Bianchi", "Luca Verdi"]
    # results = extractor.process_person_list(person_names)
    
    # CASO 3: Se lo Step 1 produce un file con i nomi
    # with open("step1_output.txt", "r") as f:
    #     person_names = [line.strip() for line in f if line.strip()]
    # results = extractor.process_person_list(person_names)
    
    # CASO 4: Con OCR se i PPTX hanno immagini
    # extractor_ocr = PPTXToJSONExtractor(
    #     cv_ppt_folder="./cv_ppt",
    #     cv_json_folder="./cv_json",
    #     use_ocr=True
    # )
    # extractor_ocr.process_person_list(person_names)


if __name__ == "__main__":
    main()